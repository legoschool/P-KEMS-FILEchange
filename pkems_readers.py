# -*- coding: utf-8 -*-
"""
PKEMS 문서 읽기 모듈
=====================
여러 형식의 문서를 '마크다운 본문'으로 읽어들인다.

지원 형식
    .hwp   한글 (구버전, HWP 5.0 바이너리)
    .hwpx  한글 (신버전, ZIP+XML)
    .docx  워드
    .pptx  파워포인트  (슬라이드별 + 발표자 노트)
    .xlsx  엑셀        (시트별 마크다운 표)
    .csv   표 데이터
    .pdf   PDF (텍스트형)
    .html  웹문서
    .txt   일반 텍스트
    .md    마크다운 (그대로 통과)
    .gdoc/.gsheet/.gslides  구글 문서 바로가기 (문서 ID만 읽음)

사용법
    from pkems_readers import read_any, SUPPORTED
    doc = read_any("보고서.hwp")
    print(doc.text)

각 읽기 함수는 ReadResult 를 돌려준다. 실패해도 예외를 던지지 않고
ok=False 와 error 메시지를 담아 돌려주므로, 일괄 변환이 중단되지 않는다.

PKEMS(개인지식경험관리체계) 프로젝트
"""

from __future__ import annotations

import os
import re
import io
import csv
import json
import zlib
import struct
import zipfile
import xml.etree.ElementTree as ET
from dataclasses import dataclass, field


# ─────────────────────────────────────────────────────────────
@dataclass
class ReadResult:
    ok: bool
    text: str = ""
    kind: str = ""                      # 형식 이름 (한글, 워드 …)
    meta: dict = field(default_factory=dict)
    error: str = ""

    @property
    def chars(self) -> int:
        return len(self.text)


def _clean(paras: list[str]) -> str:
    """빈 줄 정리 후 문단 사이 한 줄 띄우기"""
    out, prev_blank = [], True
    for p in paras:
        s = (p or "").strip()
        if not s:
            prev_blank = True
            continue
        if not prev_blank and out:
            out.append("")
        out.append(s)
        prev_blank = False
    return "\n\n".join(x for x in out if x)


# ─────────────────────────────────────────────────────────────
# 한글 (.hwp) — HWP 5.0 바이너리
# ─────────────────────────────────────────────────────────────
_TAG = 0x10
HWPTAG_PARA_TEXT = _TAG + 51            # 0x43
HWPTAG_CTRL_HEADER = _TAG + 55          # 0x47
HWPTAG_LIST_HEADER = _TAG + 56          # 0x48
HWPTAG_TABLE = _TAG + 61                # 0x4D

# 표 셀 속성은 LIST_HEADER 의 8번째 바이트부터 시작한다.
#   0  INT32  문단 수
#   4  UINT32 속성
#   8  UINT16 열(col) / 10 행(row) / 12 열병합 / 14 행병합
_CELL_OFFSET = 8
_MAX_SIDE = 300        # 한 변이 이보다 크면 표로 보지 않는다
_MAX_CELLS = 20000     # 칸이 이보다 많으면 표 대신 글로 풀어쓴다

# 문단 텍스트에 섞인 제어문자: 아래 값들은 '자기 + 6워드 + 자기' = 8워드 블록
_HWP_BLOCK_CTRL = {1, 2, 3, 4, 5, 6, 7, 8, 9, 11, 12, 14, 15,
                   16, 17, 18, 19, 20, 21, 22, 23}


def _hwp_decode_para(payload: bytes) -> str:
    """HWP 문단 레코드는 UTF-16 '코드 단위' 배열이다.
    이모지 등은 서로게이트 쌍(2워드)으로 들어오므로 합쳐 주어야 한다."""
    n = len(payload) // 2
    if n == 0:
        return ""
    words = struct.unpack_from(f"<{n}H", payload, 0)
    buf, i = [], 0
    while i < n:
        c = words[i]
        if c in _HWP_BLOCK_CTRL:
            i += 8                       # 표·그림 등 제어 블록 건너뛰기
            continue
        if c < 32:
            if c in (10, 13):
                buf.append("\n")
            elif c in (24, 30, 31):
                buf.append(" ")
            i += 1
            continue
        # 서로게이트 쌍 합치기
        if 0xD800 <= c <= 0xDBFF and i + 1 < n and 0xDC00 <= words[i + 1] <= 0xDFFF:
            buf.append(chr(0x10000 + ((c - 0xD800) << 10) + (words[i + 1] - 0xDC00)))
            i += 2
            continue
        if 0xD800 <= c <= 0xDFFF:        # 짝 없는 서로게이트는 버린다
            i += 1
            continue
        buf.append(chr(c))
        i += 1
    return "".join(buf).strip()


class _HwpTable:
    """표 하나를 모아 두었다가 마크다운 표로 내놓는다."""

    def __init__(self, level: int):
        self.level = level          # 이 표를 감싼 CTRL_HEADER 의 깊이
        self.rows = self.cols = 0
        self.cells: dict[tuple[int, int], list[str]] = {}
        self.spans: dict[tuple[int, int], tuple[int, int]] = {}
        self.cur: tuple[int, int] | None = None

    def set_size(self, payload: bytes):
        if len(payload) >= 8:
            _, self.rows, self.cols = struct.unpack_from("<IHH", payload, 0)

    def start_cell(self, payload: bytes):
        """LIST_HEADER 는 표 셀 말고 글상자 등에도 쓰인다.
        표가 선언한 크기를 벗어나는 값이면 셀이 아니라고 보고 무시한다."""
        self.cur = None
        if len(payload) < _CELL_OFFSET + 8:
            return
        col, row, cspan, rspan = struct.unpack_from("<HHHH", payload, _CELL_OFFSET)
        if self.rows and self.cols:
            if row >= self.rows or col >= self.cols:
                return                      # 표 밖 -> 셀 아님
        elif row > _MAX_SIDE or col > _MAX_SIDE:
            return                          # 크기를 모를 땐 상식선에서 자름
        self.cur = (row, col)
        self.cells.setdefault(self.cur, [])
        self.spans[self.cur] = (max(cspan, 1), max(rspan, 1))

    def add_text(self, text: str) -> bool:
        if self.cur is None:
            return False
        if text.strip():
            self.cells[self.cur].append(text.strip())
        return True

    def to_markdown(self) -> str:
        if not self.cells:
            return ""
        maxr = max(r for r, _ in self.cells) + 1
        maxc = max(c for _, c in self.cells) + 1
        # 선언 크기가 있으면 그것을 믿되, 실제 셀이 더 많으면 거기까지만 늘린다
        nrows = min(max(self.rows, maxr), _MAX_SIDE)
        ncols = min(max(self.cols, maxc), _MAX_SIDE)
        if nrows < 1 or ncols < 1:
            return ""
        if nrows * ncols > _MAX_CELLS:
            # 표로 그리기엔 너무 크다 -> 내용만 줄줄이 적는다
            return "\n\n".join(" ".join(v) for v in self.cells.values() if v)

        grid = [["" for _ in range(ncols)] for _ in range(nrows)]
        for (r, c), parts in self.cells.items():
            if r < nrows and c < ncols:
                grid[r][c] = " ".join(parts).replace("|", "／")

        # 내용이 전혀 없는 표는 버린다
        if not any(any(x for x in row) for row in grid):
            return ""

        head = grid[0]
        out = ["| " + " | ".join(head) + " |",
               "|" + "|".join(["---"] * ncols) + "|"]
        for row in grid[1:]:
            out.append("| " + " | ".join(row) + " |")
        return "\n".join(out)


def read_hwp(path: str) -> ReadResult:
    try:
        import olefile
    except ImportError:
        return ReadResult(False, kind="한글", error="olefile 설치 필요 (pip install olefile)")
    # .hwp 인데 속은 다른 형식인 경우가 있다 (hwpx 를 이름만 바꿨거나, 아주 옛 버전)
    try:
        with open(path, "rb") as fh:
            head = fh.read(8)
    except OSError as e:
        return ReadResult(False, kind="한글", error=f"파일을 열지 못했습니다: {e}")

    if head.startswith(_ZIP_MAGIC):
        return read_hwpx(path)          # 사실은 hwpx 였다 — 그대로 처리해 준다
    if not head.startswith(_OLE_MAGIC):
        return ReadResult(
            False, kind="한글",
            error=("한글 5.0 이상 형식이 아닙니다. 아주 옛 한글 문서이거나 파일이 "
                   "깨졌을 수 있습니다. 한글에서 열어 '다른 이름으로 저장'으로 "
                   ".hwp 또는 .hwpx 로 다시 저장한 뒤 변환해 주세요."))

    try:
        f = olefile.OleFileIO(path)
    except Exception as e:
        return ReadResult(False, kind="한글", error=f"파일 열기 실패: {e}")

    try:
        dirs = f.listdir()
        header = f.openstream("FileHeader").read()
        compressed = bool(header[36] & 1)

        sections = sorted(
            (d for d in dirs if d and d[0] == "BodyText" and d[1].startswith("Section")),
            key=lambda d: int(re.sub(r"\D", "", d[1]) or 0),
        )
        if not sections:
            return ReadResult(False, kind="한글", error="본문(BodyText)이 없습니다")

        paras: list[str] = []
        stack: list[_HwpTable] = []     # 표 안의 표까지 다룬다
        n_tables = 0

        def close_tables(level: int):
            """깊이가 얕아지면 그 안에서 열린 표들을 끝낸다."""
            while stack and level <= stack[-1].level:
                md = stack.pop().to_markdown()
                if md:
                    (stack[-1].add_text(md) if stack else None) or paras.append(md)

        for sec in sections:
            data = f.openstream(sec).read()
            if compressed:
                try:
                    data = zlib.decompress(data, -15)
                except zlib.error:
                    continue
            i, n = 0, len(data)
            while i < n - 4:
                (word,) = struct.unpack_from("<I", data, i)
                tag = word & 0x3FF
                level = (word >> 10) & 0x3FF
                size = (word >> 20) & 0xFFF
                i += 4
                if size == 0xFFF:
                    (size,) = struct.unpack_from("<I", data, i)
                    i += 4
                payload = data[i:i + size]
                i += size

                close_tables(level)

                if tag == HWPTAG_CTRL_HEADER and payload[:4][::-1] == b"tbl ":
                    stack.append(_HwpTable(level))
                    n_tables += 1
                elif tag == HWPTAG_TABLE and stack:
                    stack[-1].set_size(payload)
                elif tag == HWPTAG_LIST_HEADER and stack:
                    stack[-1].start_cell(payload)
                elif tag == HWPTAG_PARA_TEXT:
                    text = _hwp_decode_para(payload)
                    if not (stack and stack[-1].add_text(text)):
                        paras.append(text)

            close_tables(0)

        return ReadResult(True, _clean(paras), "한글",
                          {"섹션": len(sections), "문단": len(paras), "표": n_tables})
    except Exception as e:
        return ReadResult(False, kind="한글", error=f"본문 해석 실패: {e}")
    finally:
        try:
            f.close()
        except Exception:
            pass


# ─────────────────────────────────────────────────────────────
# 한글 (.hwpx) — ZIP + XML (표준 라이브러리만 사용)
# ─────────────────────────────────────────────────────────────
def _localname(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def _hwpx_cell_text(tc) -> str:
    """표 셀 안의 글을 모은다 (셀 안의 표까지 포함)."""
    return " ".join(
        (t.text or "").strip() for t in tc.iter()
        if _localname(t.tag) == "t" and (t.text or "").strip()
    ).strip()


def _hwpx_table_md(tbl) -> str:
    """<hp:tbl> 을 마크다운 표로 옮긴다."""
    try:
        nrows = int(tbl.get("rowCnt") or 0)
        ncols = int(tbl.get("colCnt") or 0)
    except ValueError:
        nrows = ncols = 0

    cells: dict[tuple[int, int], str] = {}
    for tc in tbl.iter():
        if _localname(tc.tag) != "tc":
            continue
        addr = next((a for a in tc if _localname(a.tag) == "cellAddr"), None)
        if addr is None:
            continue
        try:
            c = int(addr.get("colAddr", 0))
            r = int(addr.get("rowAddr", 0))
        except ValueError:
            continue
        if r > _MAX_SIDE or c > _MAX_SIDE:
            continue
        cells[(r, c)] = _hwpx_cell_text(tc)

    if not cells:
        return ""
    nrows = min(max(nrows, max(r for r, _ in cells) + 1), _MAX_SIDE)
    ncols = min(max(ncols, max(c for _, c in cells) + 1), _MAX_SIDE)
    if nrows * ncols > _MAX_CELLS:
        return "\n\n".join(v for v in cells.values() if v)
    if not any(cells.values()):
        return ""

    grid = [["" for _ in range(ncols)] for _ in range(nrows)]
    for (r, c), v in cells.items():
        if r < nrows and c < ncols:
            grid[r][c] = v.replace("|", "／")
    out = ["| " + " | ".join(grid[0]) + " |",
           "|" + "|".join(["---"] * ncols) + "|"]
    for row in grid[1:]:
        out.append("| " + " | ".join(row) + " |")
    return "\n".join(out)


def read_hwpx(path: str) -> ReadResult:
    try:
        z = zipfile.ZipFile(path)
    except Exception as e:
        return ReadResult(False, kind="한글", error=f"파일 열기 실패: {e}")
    try:
        secs = sorted(n for n in z.namelist()
                      if re.match(r"Contents/section\d+\.xml$", n))
        if not secs:
            return ReadResult(False, kind="한글", error="section XML을 찾지 못했습니다")

        paras: list[str] = []
        n_tables = 0

        def walk(el, buf: list[str]):
            """문서 순서대로 훑되, 표를 만나면 통째로 옮기고 더 내려가지 않는다."""
            nonlocal n_tables
            name = _localname(el.tag)
            if name == "tbl":
                if buf:
                    paras.append(" ".join(buf).strip())
                    buf.clear()
                md = _hwpx_table_md(el)
                if md:
                    paras.append(md)
                n_tables += 1
                return
            if name == "t" and (el.text or "").strip():
                buf.append(el.text.strip())
            for ch in el:
                walk(ch, buf)
            if name == "p" and buf:
                paras.append(" ".join(buf).strip())
                buf.clear()

        for s in secs:
            root = ET.fromstring(z.read(s))
            leftover: list[str] = []
            walk(root, leftover)
            if leftover:
                paras.append(" ".join(leftover).strip())

        return ReadResult(True, _clean(paras), "한글",
                          {"섹션": len(secs), "문단": len(paras), "표": n_tables})
    except Exception as e:
        return ReadResult(False, kind="한글", error=f"본문 해석 실패: {e}")
    finally:
        z.close()


# ─────────────────────────────────────────────────────────────
# 워드 (.docx)
# ─────────────────────────────────────────────────────────────
def read_docx(path: str) -> ReadResult:
    bad = _check_ooxml(path, "워드", ".docx")
    if bad:
        return bad
    try:
        import docx
    except ImportError:
        return ReadResult(False, kind="워드", error="python-docx 설치 필요")
    try:
        d = docx.Document(path)
        out = []
        for p in d.paragraphs:
            s = p.text.strip()
            if not s:
                continue
            style = (p.style.name or "").lower()
            m = re.search(r"heading (\d)", style)
            out.append(("#" * min(int(m.group(1)), 6) + " " + s) if m else s)
        for t in d.tables:
            out.append(_rows_to_md([[c.text.strip() for c in r.cells] for r in t.rows]))
        return ReadResult(True, _clean(out), "워드",
                          {"문단": len(d.paragraphs), "표": len(d.tables)})
    except Exception as e:
        return ReadResult(False, kind="워드", error=str(e))


# ─────────────────────────────────────────────────────────────
# 파워포인트 (.pptx)
# ─────────────────────────────────────────────────────────────
def read_pptx(path: str) -> ReadResult:
    bad = _check_ooxml(path, "파워포인트", ".pptx")
    if bad:
        return bad
    try:
        from pptx import Presentation
    except ImportError:
        return ReadResult(False, kind="파워포인트", error="python-pptx 설치 필요")
    try:
        prs = Presentation(path)
        out = []
        for i, slide in enumerate(prs.slides, 1):
            out.append(f"## 슬라이드 {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        s = "".join(r.text for r in para.runs).strip()
                        if s:
                            out.append(s)
                if getattr(shape, "has_table", False):
                    out.append(_rows_to_md(
                        [[c.text.strip() for c in r.cells] for r in shape.table.rows]))
            try:
                if slide.has_notes_slide:
                    note = slide.notes_slide.notes_text_frame.text.strip()
                    if note:
                        out += ["> **발표자 노트**", "> " + note.replace("\n", "\n> ")]
            except Exception:
                pass
        return ReadResult(True, _clean(out), "파워포인트",
                          {"슬라이드": len(prs.slides)})
    except Exception as e:
        return ReadResult(False, kind="파워포인트", error=str(e))


# ─────────────────────────────────────────────────────────────
# 엑셀 (.xlsx) / csv
# ─────────────────────────────────────────────────────────────
def _rows_to_md(rows: list[list[str]], max_rows: int = 300) -> str:
    rows = [r for r in rows if any((c or "").strip() for c in r)]
    if not rows:
        return ""
    cut = rows[:max_rows]
    width = max(len(r) for r in cut)
    def fix(r):
        r = list(r) + [""] * (width - len(r))
        return [str(c or "").replace("|", "／").replace("\n", " ").strip() for c in r]
    head = fix(cut[0])
    lines = ["| " + " | ".join(head) + " |",
             "|" + "|".join(["---"] * width) + "|"]
    for r in cut[1:]:
        lines.append("| " + " | ".join(fix(r)) + " |")
    if len(rows) > max_rows:
        lines.append(f"\n*(전체 {len(rows)}행 중 {max_rows}행만 표시)*")
    return "\n".join(lines)


def read_xlsx(path: str) -> ReadResult:
    bad = _check_ooxml(path, "엑셀", ".xlsx")
    if bad:
        return bad
    try:
        import openpyxl
    except ImportError:
        return ReadResult(False, kind="엑셀", error="openpyxl 설치 필요")
    try:
        wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
        out = []
        for ws in wb.worksheets:
            rows = [[("" if c is None else str(c)) for c in row]
                    for row in ws.iter_rows(values_only=True)]
            table = _rows_to_md(rows)
            if table:
                out += [f"## {ws.title}", table]
        wb.close()
        return ReadResult(True, _clean(out), "엑셀", {"시트": len(wb.worksheets)})
    except Exception as e:
        return ReadResult(False, kind="엑셀", error=str(e))


def read_csv(path: str) -> ReadResult:
    for enc in ("utf-8-sig", "cp949", "utf-8"):
        try:
            with io.open(path, encoding=enc, newline="") as f:
                rows = list(csv.reader(f))
            return ReadResult(True, _rows_to_md(rows), "표", {"행": len(rows)})
        except UnicodeDecodeError:
            continue
        except Exception as e:
            return ReadResult(False, kind="표", error=str(e))
    return ReadResult(False, kind="표", error="문자 인코딩을 알 수 없습니다")


# ─────────────────────────────────────────────────────────────
# HTML / 텍스트
# ─────────────────────────────────────────────────────────────
def read_html(path: str) -> ReadResult:
    raw = None
    for enc in ("utf-8", "cp949", "euc-kr"):
        try:
            with io.open(path, encoding=enc) as f:
                raw = f.read()
            break
        except (UnicodeDecodeError, LookupError):
            continue
    if raw is None:
        return ReadResult(False, kind="웹문서", error="문자 인코딩을 알 수 없습니다")
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(raw, "html.parser")
        for t in soup(["script", "style", "noscript"]):
            t.decompose()
        title = (soup.title.string or "").strip() if soup.title else ""
        parts = []
        for el in soup.find_all(["h1", "h2", "h3", "h4", "p", "li", "td", "th"]):
            s = el.get_text(" ", strip=True)
            if not s:
                continue
            if el.name.startswith("h"):
                parts.append("#" * int(el.name[1]) + " " + s)
            elif el.name == "li":
                parts.append("- " + s)
            else:
                parts.append(s)
        return ReadResult(True, _clean(parts), "웹문서", {"제목": title})
    except ImportError:
        text = re.sub(r"<[^>]+>", " ", raw)
        return ReadResult(True, _clean(text.split("\n")), "웹문서", {})
    except Exception as e:
        return ReadResult(False, kind="웹문서", error=str(e))


def read_text(path: str) -> ReadResult:
    for enc in ("utf-8", "cp949", "euc-kr", "utf-16"):
        try:
            with io.open(path, encoding=enc) as f:
                return ReadResult(True, f.read().strip(), "텍스트", {"인코딩": enc})
        except (UnicodeDecodeError, LookupError):
            continue
        except Exception as e:
            return ReadResult(False, kind="텍스트", error=str(e))
    return ReadResult(False, kind="텍스트", error="문자 인코딩을 알 수 없습니다")


# ─────────────────────────────────────────────────────────────
# PDF (일반 문서용 · 블로그 백업은 pkems_converter 를 사용)
# ─────────────────────────────────────────────────────────────
def read_pdf(path: str) -> ReadResult:
    try:
        import pymupdf
    except ImportError:
        try:
            import fitz as pymupdf
        except ImportError:
            return ReadResult(False, kind="PDF", error="pymupdf 설치 필요")
    try:
        doc = pymupdf.open(path)
        pages = doc.page_count
        out = []
        for i in range(pages):
            t = doc[i].get_text().strip()
            if t:
                out.append(t)
        doc.close()
        text = _clean(out)
        if len(text) < 20 and pages > 0:
            # 종이를 찍어 만든 PDF 는 글자가 아니라 그림이라 뽑아낼 것이 없다.
            # 글자 인식(OCR)은 이 도구의 범위를 벗어나므로 분명히 알린다.
            return ReadResult(
                False, kind="PDF",
                error=(f"글자가 없는 PDF 입니다({pages}쪽). 스캔하거나 사진으로 만든 "
                       f"문서로 보입니다. 이 도구는 글자 인식(OCR)을 하지 않으므로 "
                       f"변환할 수 없습니다."))
        return ReadResult(True, text, "PDF", {"쪽": pages})
    except Exception as e:
        return ReadResult(False, kind="PDF", error=str(e))


# ─────────────────────────────────────────────────────────────
# 구글 문서 바로가기 (.gdoc / .gsheet / .gslides)
# ─────────────────────────────────────────────────────────────
# ─────────────────────────────────────────────────────────────
# 구형 오피스 (.xls / .ppt / .doc) — 다루지 않고, 어떻게 하면 되는지 알려준다
# ─────────────────────────────────────────────────────────────
_OLD_OFFICE = {
    ".xls": ("엑셀", ".xlsx"),
    ".ppt": ("파워포인트", ".pptx"),
    ".doc": ("워드", ".docx"),
}


_OLE_MAGIC = b"\xd0\xcf\x11\xe0"       # 옛 오피스·한글의 CFB 서명
_ZIP_MAGIC = b"PK"                     # docx·pptx·xlsx·hwpx 는 모두 ZIP


def _check_ooxml(path: str, kind: str, newext: str) -> ReadResult | None:
    """확장자만 새 형식으로 바꿔 놓은 파일을 알아본다.
    맞으면 None, 아니면 안내가 담긴 ReadResult 를 돌려준다."""
    try:
        with open(path, "rb") as f:
            head = f.read(4)
    except OSError as e:
        return ReadResult(False, kind=kind, error=f"파일을 열지 못했습니다: {e}")
    if head.startswith(_ZIP_MAGIC):
        return None
    if head.startswith(_OLE_MAGIC):
        return ReadResult(
            False, kind=kind,
            error=(f"이름만 {newext} 이고 실제로는 옛 형식인 파일입니다. "
                   f"해당 파일을 열어 '다른 이름으로 저장'으로 진짜 {newext} 형식으로 "
                   f"바꾼 뒤 다시 변환해 주세요."))
    return ReadResult(False, kind=kind, error=f"{newext} 형식이 아닙니다 (내용이 깨졌을 수 있음)")


def read_old_office(path: str) -> ReadResult:
    """옛 형식은 구조가 완전히 달라 따로 다루지 않는다.
    해당 프로그램에서 '다른 이름으로 저장'만 하면 되므로 그 방법을 알려준다."""
    ext = os.path.splitext(path)[1].lower()
    kind, newext = _OLD_OFFICE.get(ext, ("문서", ".xlsx"))
    return ReadResult(
        False, kind=kind,
        error=(f"옛 {kind} 형식({ext})은 지원하지 않습니다. "
               f"해당 파일을 열어 '다른 이름으로 저장'으로 {newext} 형식으로 "
               f"바꾼 뒤 다시 변환해 주세요."))


_G_KIND = {".gdoc": "구글문서", ".gsheet": "구글시트", ".gslides": "구글슬라이드"}


def read_gshortcut(path: str) -> ReadResult:
    """구글 드라이브 바로가기 파일에서 문서 ID/주소만 읽는다.

    구글 문서·시트·슬라이드는 '내 컴퓨터에 실체가 없는' 온라인 문서다.
    윈도우 드라이브 앱에서는 파일로 열리지 않는 경우가 많으므로(가상 파일),
    실제 내용은 Drive API 로 내보내야 한다(pkems_gdrive.export_google_doc).
    """
    ext = os.path.splitext(path)[1].lower()
    kind = _G_KIND.get(ext, "구글문서")
    guide = (f"> 구글 {kind}입니다. 내용이 온라인에만 있어 파일로는 읽을 수 없습니다.\n"
             f"> 코랩에서 'Drive API 내보내기'로 가져와야 합니다.")
    try:
        with io.open(path, encoding="utf-8") as f:
            info = json.load(f)
        doc_id = info.get("doc_id") or info.get("resource_id", "").split(":")[-1]
        url = info.get("url", "")
        return ReadResult(True, f"{guide}\n\n{url}", kind,
                          {"doc_id": doc_id, "url": url, "needs_api": True})
    except OSError:
        # 드라이브 앱의 가상 파일 — 열람 자체가 불가
        return ReadResult(True, guide, kind,
                          {"needs_api": True, "note": "가상 파일이라 로컬에서 열 수 없음"})
    except Exception as e:
        return ReadResult(False, kind=kind, error=str(e))


# ─────────────────────────────────────────────────────────────
# 등록표
# ─────────────────────────────────────────────────────────────
READERS = {
    ".hwp": read_hwp,
    ".hwpx": read_hwpx,
    ".docx": read_docx,
    ".pptx": read_pptx,
    ".xlsx": read_xlsx,
    ".xlsm": read_xlsx,
    ".csv": read_csv,
    ".tsv": read_csv,
    ".pdf": read_pdf,
    ".html": read_html,
    ".htm": read_html,
    ".txt": read_text,
    ".md": read_text,
    ".gdoc": read_gshortcut,
    ".gsheet": read_gshortcut,
    ".gslides": read_gshortcut,
    # 옛 형식 — 변환하지 않고 '어떻게 바꾸면 되는지' 안내만 남긴다
    ".xls": read_old_office,
    ".ppt": read_old_office,
    ".doc": read_old_office,
}

SUPPORTED = sorted(READERS)


def sanitize(text: str) -> str:
    """파일로 저장할 수 없는 글자(짝 없는 서로게이트 등)를 걸러낸다."""
    if not text:
        return text
    try:
        text.encode("utf-8")
        return text
    except UnicodeEncodeError:
        return text.encode("utf-8", "ignore").decode("utf-8", "ignore")


def read_any(path: str) -> ReadResult:
    """확장자를 보고 알맞은 읽기 함수를 고른다."""
    ext = os.path.splitext(path)[1].lower()
    fn = READERS.get(ext)
    if fn is None:
        return ReadResult(False, kind=ext or "?", error="지원하지 않는 형식")
    if not os.path.exists(path):
        return ReadResult(False, kind=ext, error="파일이 없습니다")
    try:
        res = fn(path)
    except Exception as e:                       # 어떤 경우에도 죽지 않게
        return ReadResult(False, kind=ext, error=f"예기치 못한 오류: {e}")
    res.text = sanitize(res.text)
    return res
